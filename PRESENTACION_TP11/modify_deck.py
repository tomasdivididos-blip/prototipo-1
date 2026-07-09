# -*- coding: utf-8 -*-
"""Modifica el deck del usuario (TP11_Presentacion.pptx): imagen slide 3,
slide de scores, slide de respuestas, y arregla la numeracion."""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__)); A = os.path.join(HERE, 'assets')
F = os.path.join(HERE, 'TP11_Presentacion.pptx')
F_IN = os.path.join(HERE, 'TP11_Presentacion_BACKUP.pptx')   # leer del backup limpio
BLUE='2b6cb0'; DARK='202020'; GREY='555555'; FOOT='Diseño modal de recintos pequeños'
prs = Presentation(F_IN)
slides = list(prs.slides)
print('slides iniciales:', len(slides))

def tb(s,l,t,w,h,text,size=14,bold=False,color=DARK,align=PP_ALIGN.LEFT):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=b.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.name='Arial'; r.font.color.rgb=RGBColor.from_string(color); return b
def bullets(s,l,t,w,h,items,size=15,color=DARK,gap=6):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=b.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        r=p.add_run(); r.text='•  '+it; r.font.size=Pt(size); r.font.name='Arial'; r.font.color.rgb=RGBColor.from_string(color)
    return b
def img(s,name,l,t,w=None,h=None):
    kw={}
    if w:kw['width']=Inches(w)
    if h:kw['height']=Inches(h)
    return s.shapes.add_picture(os.path.join(A,name),Inches(l),Inches(t),**kw)
def set_ph(s,idx,text,size=None,bold=None):
    for ph in s.placeholders:
        if ph.placeholder_format.idx==idx:
            ph.text=text
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    if size:r.font.size=Pt(size)
                    if bold is not None:r.font.bold=bold
            return
def footer_boxes(s):
    tb(s,3.25,6.95,3.4,0.4,FOOT,10,False,GREY,PP_ALIGN.CENTER)
    tb(s,7.55,6.95,1.75,0.4,'',10,False,GREY,PP_ALIGN.RIGHT)   # numero -> lo llena el renumber

# =============== 1) IMAGEN SLIDE 3 (indice 2) ===============
s3 = slides[2]
pic = None
for sh in s3.shapes:
    if sh.shape_type == 13:   # PICTURE
        pic = sh
if pic is None:
    raise RuntimeError('no encontre la imagen en slide 3')
L,T,W = pic.left, pic.top, pic.width
pic._element.getparent().remove(pic._element)
s3.shapes.add_picture(os.path.join(A,'fig_mode111.png'), L, T, width=W)
# actualizar caption (1,1,0) -> (1,1,1)
for sh in s3.shapes:
    if sh.has_text_frame and '1,1,0' in sh.text_frame.text:
        sh.text_frame.text = 'Forma modal (1,1,1): planos nodales en el centro (x=0 y z=1,5) y antinodos en las esquinas.'
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.font.size=Pt(10); r.font.name='Arial'; r.font.color.rgb=RGBColor.from_string(GREY)
print('slide 3: imagen reemplazada por (1,1,1)')

# =============== 2) SLIDE DE SCORES (despues de "tres ejes") ===============
sc = prs.slides.add_slide(prs.slide_layouts[1])
set_ph(sc,0,'¿Cómo se ponderan los scores?',22,True); footer_boxes(sc)
tb(sc,0.72,1.6,8.6,0.3,'Grupo MODAL — peso interno de cada sub-score',12,True,BLUE)
tb(sc,0.9,1.95,8.4,0.7,'Bolt-spacing 0,25  ·  FSI 0,15  ·  Bonello 0,05  ·  Modal-Q/Fazenda 0,20  ·  RT60 0,20  ·  Schroeder 0,15',14)
bullets(sc,0.72,2.7,8.6,1.0,[
 'Las métricas baratas y que más discriminan (Bolt, aspecto) pesan más; las dominadas por V o RT60 (inputs constantes), menos.'],14)
tb(sc,0.72,3.75,8.6,0.3,'Peso de cada categoría en el score total, según el uso',12,True,BLUE)
img(sc,'tabla3b_pesos.png',0.7,4.1,w=6.6)
bullets(sc,0.72,6.05,8.6,0.6,[
 'Score total = suma ponderada de las 5 categorías. En el eje Ubicación los pesos son ajustables por el usuario.'],13)
# mover despues de "tres ejes"
ids=list(prs.slides._sldIdLst); el=ids[-1]
tres_idx=None
for i,sl in enumerate(prs.slides):
    for sh in sl.shapes:
        try:
            if sh.placeholder_format.idx==0 and 'tres ejes' in sh.text_frame.text.lower(): tres_idx=i
        except: pass
prs.slides._sldIdLst.remove(el)
prs.slides._sldIdLst.insert(tres_idx+1, el)
print('slide de scores insertada tras "tres ejes" (idx %s)'%tres_idx)

# =============== 3) SLIDE DE RESPUESTAS (al final) ===============
ans = prs.slides.add_slide(prs.slide_layouts[1])
set_ph(ans,0,'Respuestas anticipadas (backup)',22,True); footer_boxes(ans)
bullets(ans,0.72,1.7,8.6,5.0,[
 '¿Por qué P1 y no P2?  P2 baja el error a <0,05 % pero cuesta 5–36× más; P1 (0,4–3 %) ya queda por debajo del ruido del modelado físico. Mejor compromiso.',
 '¿Por qué no un algoritmo genético?  En geometría el espacio está dominado por ratios óptimos conocidos (sembrarlos es más barato y reproducible); en ubicación cada evaluación cuesta un FEM. El AG sí aplica al espacio de quiebres (trabajo futuro).',
 '¿Malla escalonada vs boundary-fitted?  Con paredes rígidas la condición de Neumann es natural en la forma débil; la frontera escalonada no la compromete. Error volumétrico 1–2 % en los primeros modos.',
 '¿Damping modal vs matriz de impedancia?  La impedancia rompe la simetría de K y M y obliga a un solver complejo, más lento. El damping modal desde RT60 concuerda dentro de ~2 dB en la banda modal.'], size=14, gap=12)
print('slide de respuestas agregada al final')

# =============== 4) RENUMERAR ===============
allsl=list(prs.slides); total=len(allsl)
pat=re.compile(r'^\s*\d+\s*/\s*\d+\s*$')
fixed=0
for pos,sl in enumerate(allsl,1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            tx=sh.text_frame.text.strip()
            is_num = (pat.match(tx) or 'TotNum' in tx or
                      (tx=='' and sh.left is not None and abs(Emu(sh.left).inches-7.55)<0.15 and abs(Emu(sh.top).inches-6.95)<0.15))
            if is_num:
                sh.text_frame.text='%d / %d'%(pos,total)
                for p in sh.text_frame.paragraphs:
                    p.alignment=PP_ALIGN.RIGHT
                    for r in p.runs: r.font.size=Pt(10); r.font.name='Arial'; r.font.color.rgb=RGBColor.from_string(GREY)
                fixed+=1
print('numeros corregidos:', fixed, '| total slides:', total)

prs.save(F)
print('GUARDADO ->', F)
