# -*- coding: utf-8 -*-
"""Construye el deck completo (26 slides) sobre el template UNTREF."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

HERE = os.path.dirname(os.path.abspath(__file__)); A = os.path.join(HERE, 'assets')
TPL = r'C:\Users\aceve\Downloads\Template_Slide-UNTREF_OFICIAL.pptx'
BLUE = '2b6cb0'; ORANGE = 'c05621'; DARK = '202020'; GREY = '555555'
FOOT = 'Diseño modal de recintos pequeños'; NTOT = 26

prs = Presentation(TPL)
xml = prs.slides._sldIdLst
for s in list(xml):
    prs.part.drop_rel(s.get(qn('r:id'))); xml.remove(s)

# editar textos de layout (una vez): pie nombre (OBJECT) + "Gracias" (seccion)
for sh in prs.slide_layouts[1].shapes:
    if sh.has_text_frame and sh.left is not None and abs(Emu(sh.left).inches-0.69)<0.12 and abs(Emu(sh.top).inches-6.95)<0.12:
        sh.text_frame.text='Acevedo, Tomás'
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.font.size=Pt(10); r.font.name='Arial'
for sh in prs.slide_layouts[2].shapes:
    if sh.has_text_frame and 'Gracias' in sh.text_frame.text:
        sh.text_frame.text='Preguntas'
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.font.size=Pt(30); r.font.bold=True; r.font.name='Arial'

def L(i): return prs.slide_layouts[i]
def set_ph(s, idx, text, size=None, bold=None, color=None):
    for ph in s.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    if size: r.font.size = Pt(size)
                    if bold is not None: r.font.bold = bold
                    if color: r.font.color.rgb = RGBColor.from_string(color)
            return
def tb(s, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = 'Arial'
    r.font.color.rgb = RGBColor.from_string(color); return b
def bullets(s, l, t, w, h, items, size=15, color=DARK, gap=6):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = '•  ' + it
        r.font.size = Pt(size); r.font.name = 'Arial'; r.font.color.rgb = RGBColor.from_string(color)
    return b
def img(s, name, l, t, w=None, h=None):
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    return s.shapes.add_picture(os.path.join(A, name), Inches(l), Inches(t), **kw)
def lab(s, l, t, w, text, size=13):
    return tb(s, l, t, w, 0.32, text, size, True, BLUE)

_n = [0]
def content(title, tsize=22):
    _n[0] += 1
    s = prs.slides.add_slide(L(1))
    set_ph(s, 0, title, tsize, True)
    tb(s, 3.25, 6.95, 3.4, 0.4, FOOT, 10, False, GREY, PP_ALIGN.CENTER)
    tb(s, 7.55, 6.95, 1.75, 0.4, '%d / %d' % (_n[0], NTOT), 10, False, GREY, PP_ALIGN.RIGHT)
    return s

# ============================== S1 TÍTULO ==============================
s = prs.slides.add_slide(L(0))
set_ph(s, 0, 'Software para el diseño modal de recintos pequeños', 28, True)
tb(s, 0.85, 3.95, 8.5, 1.0, 'Predicción por elementos finitos, figuras de mérito modales y optimización de forma y fuentes', 15)
tb(s, 0.85, 5.25, 8.5, 0.4, 'Acevedo, Tomás  ·  Legajo 51323', 15, True)
tb(s, 0.85, 5.75, 8.5, 0.4, 'UNTREF · Ingeniería de Sonido · TP 11', 13, color=GREY)

# ============================== S2 PROBLEMA ==============================
s = content('El problema: acústica modal en salas chicas')
bullets(s, 0.72, 1.7, 4.7, 5.0, [
 'Por debajo de la frecuencia de Schroeder el campo es resonante, no difuso.',
 'Cada modo concentra energía en frecuencias y posiciones fijas.',
 'Resultado: coloración, refuerzos y cancelaciones en el punto de escucha.',
 'En salas de control, estudios y home theaters, esto domina la baja frecuencia.',
 'Se mitiga con la forma/proporciones y con la ubicación de fuente y oyente.'], size=15)
img(s, 'fig5_heatmap_spl.png', 5.55, 1.9, w=3.9)
tb(s, 5.55, 5.55, 3.9, 0.6, 'Modo tangencial (1,1,0): nodos oscuros, antinodos en las esquinas.', 10, False, GREY)

# ============================== S3 CONSIGNA ==============================
s = content('La consigna — TP 11')
bullets(s, 0.72, 1.7, 8.6, 3.4, [
 'Software de código abierto con GUI para hallar la mejor forma, dimensiones y ubicación de fuentes, partiendo de un paralelepípedo con quiebres de la envolvente.',
 'Herramientas sugeridas: algoritmos genéticos, figuras de mérito modales y la ecuación de onda para la respuesta modal.',
 'Entregable adicional: una pregunta de investigación sin respuesta directa en la bibliografía, abordable en 6 meses.'], size=15, gap=8)
tb(s, 0.72, 5.0, 8.6, 0.32, 'Rúbrica de evaluación (100 pts):', 13, True, BLUE)
tb(s, 0.72, 5.35, 8.6, 0.8, 'Desarrollo del código 40  ·  GUI 15  ·  Cálculo de parámetros 15  ·  Estado del arte 10  ·  Fuentes de conflicto 10  ·  Introducción + Conclusiones 10', 14)

# ============================== S4 CÓMO LO RESOLVÍ ==============================
s = content('Cómo lo resolví')
bullets(s, 0.72, 1.7, 4.7, 5.0, [
 'Predictor modal por FEM: resuelve la ecuación de onda para cada alternativa.',
 'Evaluación con figuras de mérito modales.',
 'Tres ejes de optimización: forma, ubicación de fuentes y combinado.',
 'GUI en tres pestañas; geometría con quiebres (lofting).',
 'Búsqueda determinística por semillas + refinamiento (en vez de AG).'], size=15)
img(s, 'fig1_arquitectura.png', 5.35, 2.3, w=4.1)

# ============================== S5 ONDA -> HELMHOLTZ ==============================
s = content('De la ecuación de onda a Helmholtz')
lab(s, 0.72, 1.7, 8.6, 'Ecuación de onda (presión sonora)'); img(s, 'eq_wave.png', 0.9, 2.05, w=3.6)
lab(s, 0.72, 3.5, 8.6, 'Régimen armónico → ecuación de Helmholtz'); img(s, 'eq_helmholtz.png', 0.9, 3.85, w=5.2)
tb(s, 0.72, 5.3, 8.6, 1.0, 'Buscamos las soluciones no triviales con paredes rígidas: los modos propios del recinto.', 15)

# ============================== S6 MODOS PROPIOS ==============================
s = content('Modos propios del recinto')
lab(s, 0.72, 1.7, 8.6, 'Problema de autovalores — paredes rígidas (Neumann)'); img(s, 'eq_eigen.png', 0.9, 2.05, w=6.4)
lab(s, 0.72, 3.7, 8.6, 'Solución analítica exacta (caja rectangular) = banco de prueba'); img(s, 'eq_modes.png', 0.9, 4.05, w=7.4)
tb(s, 0.72, 5.7, 8.6, 0.8, 'La caja con solución cerrada permite validar cualquier método numérico.', 15)

# ============================== S7 SCHROEDER + WEYL ==============================
s = content('Régimen modal: Schroeder y densidad')
lab(s, 0.72, 1.7, 8.6, 'Frecuencia de Schroeder (techo del régimen modal)'); img(s, 'eq_schroeder.png', 0.9, 2.05, w=3.2)
lab(s, 0.72, 3.5, 8.6, 'Densidad modal acumulada (ley de Weyl)'); img(s, 'eq_weyl.png', 0.9, 3.85, w=6.6)
tb(s, 0.72, 5.4, 8.6, 1.0, 'Por debajo de f_S los modos son discretos y el FEM es exacto; por encima, el campo es estadístico.', 15)

# ============================== S8 RESPUESTA FORZADA ==============================
s = content('Respuesta forzada: FRF modal')
lab(s, 0.72, 1.7, 8.6, 'Función de Green modal — FRF en el receptor (factor c²)'); img(s, 'eq_green.png', 0.9, 2.05, w=7.2)
lab(s, 0.72, 3.9, 8.6, 'Amortiguamiento de cada modo, derivado del RT60'); img(s, 'eq_xi.png', 0.9, 4.3, w=3.4)
tb(s, 0.72, 5.5, 8.6, 1.0, 'El factor c² sale de la derivación canónica; ξ_n controla el ancho de los picos.', 15)

# ============================== S9 RT60 ==============================
s = content('Absorción y tiempo de reverberación')
lab(s, 0.72, 1.7, 8.6, 'Sabine (absorción baja)'); img(s, 'eq_sabine.png', 0.9, 2.05, w=2.8)
lab(s, 0.72, 3.5, 8.6, 'Norris–Eyring (también válido para absorción alta)'); img(s, 'eq_eyring.png', 0.9, 3.85, w=5.6)
tb(s, 0.72, 5.4, 8.6, 1.0, 'El RT60 por banda fija el amortiguamiento modal ξ_n de la simulación.', 15)

# ============================== S10 POR QUÉ MOTOR PROPIO ==============================
s = content('¿Por qué un motor FEM y de mallado propio?')
bullets(s, 0.72, 1.7, 8.6, 5.0, [
 'Problema chico (banda modal): no justifica FEniCS/PETSc ni dependencias C++ en Windows/Anaconda.',
 'Transparencia: cada paso del ensamblaje es visible, comentado y validado con oráculos.',
 'Paredes rígidas → la condición de Neumann es natural en la forma débil: la frontera escalonada (voxel) no afecta el ensamblaje.',
 'Elementos P1 + masa consistente: error 0,4–3 % vs analítico, por debajo del ruido del modelado físico.'], size=16, gap=10)

# ============================== S11 FLOW MALLADO ==============================
s = content('Motor de mallado: cómo funciona')
img(s, 'flow_mesh.png', 0.55, 2.05, w=8.5)
tb(s, 0.72, 4.9, 8.6, 1.4, 'Voxelización + partición de Freudenthal, sin TetGen/CGAL. La frontera escalonada no compromete el ensamblaje: con paredes rígidas la condición de Neumann es natural en la forma débil.', 15)

# ============================== S12 FLOW FEM ==============================
s = content('Motor FEM: cómo funciona')
img(s, 'flow_fem.png', 0.55, 1.95, w=8.5)
lab(s, 0.72, 4.4, 8.6, 'Problema de autovalores generalizado (Lanczos + shift-invert)'); img(s, 'eq_eigprob.png', 0.9, 4.75, w=5.2)
tb(s, 0.72, 6.0, 8.6, 0.6, 'K y M reales y simétricas (paredes rígidas) → eigsh directo, eficiente.', 14)

# ============================== S13 COMPARACIÓN + CÓDIGO ==============================
s = content('Los dos motores, comparados')
tb(s, 0.72, 1.55, 8.6, 0.3, 'Mallado', 12, True, BLUE); img(s, 'flow_mesh.png', 0.6, 1.85, w=8.5)
tb(s, 0.72, 3.15, 8.6, 0.3, 'FEM', 12, True, ORANGE); img(s, 'flow_fem.png', 0.6, 3.45, w=8.5)
img(s, 'code_assembly.png', 0.9, 4.95, w=7.4)

# ============================== S14 CRITERIOS: FUENTES ==============================
s = content('¿De dónde salen los criterios de predicción?')
img(s, 'table_sources.png', 0.75, 1.75, w=8.6)

# ============================== S15 PROPORCIONES ==============================
s = content('Proporciones óptimas del recinto')
img(s, 'fig_bolt.png', 0.6, 1.8, w=4.9)
bullets(s, 5.7, 2.0, 3.7, 4.5, [
 'Cinco ternas de la literatura: Louden, Bolt, Sepmeyer, Cox y BBC/Rindel.',
 'Cada una se escala al volumen objetivo del uso.',
 'Se confinan a un rango de altura constructivo.',
 'Se puntúan por figura de mérito y se muestran las 3 mejores.'], size=14)

# ============================== S16 BONELLO + FSI ==============================
s = content('Distribución modal pareja: Bonello y FSI')
img(s, 'fig_bonello.png', 0.9, 1.8, w=7.4)
bullets(s, 0.72, 5.2, 8.6, 1.4, [
 'Bonello (1981): la densidad modal por 1/3 de octava debe ser monótona no decreciente.',
 'FSI (Rindel): mide la regularidad del espaciado entre modos consecutivos.'], size=14, gap=4)

# ============================== S17 FAZENDA ==============================
s = content('Umbral perceptual: Fazenda')
img(s, 'fig_fazenda.png', 0.75, 1.8, w=5.6)
bullets(s, 6.5, 2.2, 2.9, 4.0, [
 'Dos curvas de audibilidad modal.',
 '"Artificial": peor caso, sin enmascaramiento.',
 '"Música": escucha real.',
 'La curva la elige el programa según el uso de la sala.'], size=13)

# ============================== S18 GESTIÓN ACTIVA ==============================
s = content('Gestión activa vs. proporción')
bullets(s, 0.72, 1.7, 8.6, 2.4, [
 'Toole y Geddes: en la práctica, la ubicación de fuente/oyente, la absorción y la ecualización dominan sobre la proporción pura.',
 'Welti & Devantier: múltiples subwoofers mejoran la consistencia asiento a asiento.',
 'Interferencia fuente-frontera (SBIR): peine de realces y cancelaciones por reflexiones.'], size=15, gap=8)
lab(s, 0.72, 4.6, 8.6, 'SBIR por fuentes imagen de primer orden'); img(s, 'eq_sbir.png', 0.9, 4.95, w=7.0)

# ============================== S19 ARQUITECTURA ==============================
s = content('Arquitectura del software')
img(s, 'fig1_arquitectura.png', 1.4, 1.9, w=7.2)
tb(s, 0.72, 6.0, 8.6, 0.6, 'Núcleo de cómputo puro (numpy/scipy) desacoplado de la capa de interfaz PyQt5.', 14)

# ============================== S20 GUI ==============================
s = content('La interfaz gráfica')
img(s, 'fig3_gui.png', 3.0, 1.75, w=6.3)
bullets(s, 0.72, 2.0, 2.2, 4.5, [
 'Tres pestañas: Geometría, Acústica y Predicción.',
 'Visor 3D con manipulación directa.',
 'Fuentes (bafles) y receptor.'], size=13)

# ============================== S21 GEOMETRÍA QUIEBRES ==============================
s = content('Geometría con quiebres (lofting)')
img(s, 'fig2_editor.png', 2.1, 1.7, w=5.7)
tb(s, 0.72, 6.1, 8.6, 0.6, 'Del paralelepípedo a una envolvente quebrada: edición de planta (a) y perfiles de pared por corte lateral (b).', 13)

# ============================== S22 FIGURAS DE MÉRITO ==============================
s = content('Figuras de mérito modales')
lab(s, 0.72, 1.7, 8.6, 'Planitud media y consistencia espacial (banda válida)'); img(s, 'eq_fom.png', 0.9, 2.05, w=7.6)
bullets(s, 0.72, 3.5, 8.6, 3.0, [
 'FoM_flat: desviación de la respuesta media (planitud).',
 'FoM_esp: variación asiento a asiento (estilo Welti).',
 'Se calculan sobre una grilla de escucha, con suavizado por fracciones de octava y sólo en la banda válida de la malla.',
 'Cruce por solapamiento modal: frontera de Schroeder desde la densidad modal real.'], size=15, gap=7)

# ============================== S23 TRES EJES ==============================
s = content('Los tres ejes de predicción')
img(s, 'fig7_ejes.png', 0.55, 1.75, w=8.7)
tb(s, 0.72, 5.7, 8.6, 0.9, 'Geometría, ubicación y combinado. Búsqueda por semillas + refinamiento local (no AG): más económica, reproducible e interpretable en estos espacios.', 13)

# ============================== S24 VALIDACIÓN ==============================
s = content('Validación del solver')
img(s, 'tabla1_validacion.png', 0.55, 1.7, w=4.2)
img(s, 'fig4_frf.png', 5.1, 2.75, w=4.4)
tb(s, 0.55, 6.35, 8.8, 0.5, 'Modos analíticos vs FEM: error medio 1,6 % (máx 3,5 %). FRF con overlay de zonas no ecualizables.', 12)

# ============================== S25 PREGUNTA + CONCLUSIÓN ==============================
s = content('Pregunta de investigación y conclusión')
lab(s, 0.72, 1.7, 8.6, 'Pregunta de investigación')
tb(s, 0.72, 2.05, 8.6, 1.6, '¿Cuál es el mínimo quiebre geométrico (desplazamiento de un plano) que mejora la uniformidad modal por encima del umbral perceptible, y que supere a lo que se logra —gratis— sólo reubicando fuente y receptor?', 15)
lab(s, 0.72, 3.9, 8.6, 'Conclusión')
bullets(s, 0.72, 4.25, 8.6, 2.2, [
 'Se cumplió la consigna: predictor modal FEM + figuras de mérito + optimización en 3 ejes, con GUI.',
 'Se superó el alcance mínimo: materiales, RT60, SBIR, corregibilidad EQ, auralización, distribución Win/Mac.',
 'El espacio de quiebres queda identificado como el ámbito natural de un algoritmo genético (trabajo futuro).'], size=14, gap=6)

# ============================== S26 PREGUNTAS (sección) ==============================
s = prs.slides.add_slide(L(2))
set_ph(s, 1, 'Disparadores: ¿por qué P1 y no P2?  ·  ¿por qué no un AG?  ·  ¿malla escalonada vs boundary-fitted?  ·  ¿damping modal vs impedancia?', 15)

prs.save(os.path.join(HERE, 'TP11_Presentacion.pptx'))
print('OK slides:', len(prs.slides._sldIdLst))
